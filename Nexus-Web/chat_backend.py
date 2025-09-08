import os
import google.generativeai as genai
# REMOVIDO: import google.api_core.exceptions as api_exceptions
from pptx import Presentation

# Este es el nuevo punto de entrada de tu aplicación
def cargar_conocimiento(path):
    """
    Carga el texto de todas las formas de un archivo PowerPoint.
    """
    texto = ""
    try:
        if not os.path.exists(path):
            return "❌ Archivo 'PLAN de Capacitacion.pptx' no encontrado."

        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + " "
        return texto.strip()
    except Exception as e:
        return f"Error al cargar el archivo PowerPoint: {e}"

def consultar_gemini(pregunta, conocimiento_jira):
    """
    Genera una respuesta utilizando la API de Gemini.
    """
    try:
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            return "Error: La clave API de Gemini no está configurada. Contacta al administrador."

        genai.configure(api_key=api_key)
        model = genai.GenerativeModel("gemini-1.5-flash-latest")

        prompt = (
            f"Eres un Tester Senior con amplio conocimiento en ISTQB. Tu misión es actuar como asistente para resolver dudas de un proyecto de software, "
            f"específicamente en un contexto de pruebas de software, control de calidad y gestión de incidencias en Jira. "
            f"Debes responder a la pregunta del usuario utilizando, en primer lugar, el siguiente 'conocimiento del proyecto'. "
            f"Si el conocimiento no es suficiente, debes responder con tu conocimiento general sobre pruebas de software y control de calidad.\n\n"
            f"**Formato de Respuesta:** La respuesta debe ser facil de leer, es decir, usa saltos de linea, viñetas o cualquier otro metodo para que el parrafo generado tenga una estructura profesional y limpia\n\n"
            f"Conocimiento del proyecto:\n---\n{conocimiento_jira}\n---\n\n"
            f"Pregunta del usuario: {pregunta}"
        )

        response = model.generate_content(
            prompt,
            safety_settings=[
                {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"}
            ]
        )
        return response.text
    # CAMBIADO: Manejo genérico de excepciones en lugar de BlockedPromptException específica
    except Exception as e:
        error_message = str(e).lower()
        if "blocked" in error_message or "safety" in error_message:
            return f"Error de seguridad: La solicitud fue bloqueada por filtros de seguridad."
        else:
            return f"Error al comunicarse con la API de Gemini: {e}"
